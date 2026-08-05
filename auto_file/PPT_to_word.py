from pathlib import Path
from dataclasses import dataclass, field
import sys
from pptx import Presentation

from docx import Document
from docx.shared import Pt, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn


# =====================================================
# 配置区域
# =====================================================

if getattr(sys, "frozen", False):
    BASE_DIR = Path(sys.executable).parent
else:
    BASE_DIR = Path(__file__).resolve().parent


# 输入PPT
PPT_FILE = BASE_DIR / "input.pptx"

# 输出Word
WORD_FILE = BASE_DIR / "output.docx"

# 最大标题长度限制
MAX_TITLE_LENGTH = 30

# =====================================================
# 数据结构
# =====================================================


@dataclass
class SlideContent:
    """
    保存每一页PPT内容
    """

    page: int

    title: str = ""

    paragraphs: list[str] = field(
        default_factory=list
    )


def get_effective_font_size(
        run,
        shape
):
    """
    获取文字实际字号

    优先级:

    1. run直接设置
    2. paragraph级别
    3. shape占位符样式
    4. 默认返回0
    """

    # ======================
    # 1. run自身字号
    # ======================

    if run.font.size:

        return run.font.size.pt

    # ======================
    # 2. paragraph字号
    # ======================

    paragraph = run._parent

    if paragraph.font.size:

        return paragraph.font.size.pt

    # ======================
    # 3. placeholder样式
    # ======================

    if shape.is_placeholder:

        placeholder_format = (
            shape.placeholder_format
        )

        try:

            level = (
                paragraph.level
            )

            text_frame = (
                shape.text_frame
            )

            style = (
                text_frame._element
                .lstStyle
            )

        except Exception:

            return 0

    return 0


# =====================================================
# PPT文字提取模块
# =====================================================

def extract_ppt_text(
        ppt_file: Path
) -> list[SlideContent]:

    prs = Presentation(
        ppt_file
    )

    slides = []

    for page, slide in enumerate(
            prs.slides,
            start=1
    ):

        paragraphs = []

        title_candidate = ""

        max_font_size = 0

        for shape in slide.shapes:

            if not shape.has_text_frame:

                continue

            text = shape.text.strip()

            if not text:

                continue

            paragraphs.append(
                text
            )

            # ==========================
            # 获取当前文本框最大字号
            # ==========================

            shape_font_size = 0

            for paragraph in shape.text_frame.paragraphs:

                for run in paragraph.runs:

                    size = get_effective_font_size(
                        run,
                        shape
                    )

                    if size > shape_font_size:

                        shape_font_size = size

            # ==========================
            # 判断是否为最大字号文本框
            # ==========================

            if shape_font_size > max_font_size:

                max_font_size = shape_font_size

                title_candidate = text

        # ==========================
        # 标题有效性判断
        # ==========================

        if (
            title_candidate
            and
            len(title_candidate)
            <= MAX_TITLE_LENGTH
        ):

            title = title_candidate

            # 删除标题文本

            if title in paragraphs:

                paragraphs.remove(
                    title
                )

        else:

            title = (
                f"第{page}页"
            )

        slides.append(
            SlideContent(
                page=page,
                title=title,
                paragraphs=paragraphs
            )
        )

    return slides


# =====================================================
# 内容分析模块
# =====================================================

def analyze_content(
        slides: list[SlideContent]
) -> list[SlideContent]:
    """
    对已经完成标题识别后的内容进行整理

    功能:
    1. 保留extract_ppt_text()生成的标题
    2. 清理正文空白
    3. 删除重复文本
    4. 无标题页面生成默认标题
    """

    for slide in slides:

        # =====================================
        # 1. 检查标题
        # =====================================

        if not slide.title:

            slide.title = (
                f"第{slide.page}页"
            )

        # =====================================
        # 2. 清理正文
        # =====================================

        cleaned_paragraphs = []

        for text in slide.paragraphs:

            # 去除首尾空格

            text = text.strip()

            # 空文本跳过

            if not text:

                continue

            cleaned_paragraphs.append(
                text
            )

        slide.paragraphs = (
            cleaned_paragraphs
        )

        # =====================================
        # 3. 删除重复正文
        # =====================================

        unique_paragraphs = []

        seen = set()

        for text in slide.paragraphs:

            if text not in seen:

                unique_paragraphs.append(
                    text
                )

                seen.add(
                    text
                )

        slide.paragraphs = (
            unique_paragraphs
        )

    return slides


# =====================================================
# Word格式设置
# =====================================================

def set_word_style(
        doc: Document
):
    """
    设置Word字体格式

    正文:
        中文: 仿宋_GB2312
        英文: Times New Roman
        字号: 四号

    标题:
        中文: 方正小标宋简体
        字号: 小三
    """

    # ==========================
    # 正文 Normal 样式
    # ==========================

    normal_style = doc.styles["Normal"]

    # 英文字体

    normal_style.font.name = (
        "Times New Roman"
    )

    # 中文字体必须额外设置

    normal_style._element.rPr.rFonts.set(
        qn("w:eastAsia"),
        "仿宋_GB2312"
    )

    # 字号 四号 = 14pt

    normal_style.font.size = Pt(14)

    # ==========================
    # 标题样式
    # ==========================

    heading_style = (
        doc.styles["Heading 1"]
    )

    heading_style._element.rPr.rFonts.set(
        qn("w:eastAsia"),
        "方正小标宋简体"
    )

    # 英文标题字体

    heading_style.font.name = (
        "Times New Roman"
    )

    # 小三 = 15pt

    heading_style.font.size = Pt(15)
    # ==========================
    # 文档标题 Title 样式
    # ==========================

    title_style = doc.styles["Title"]

    title_style._element.rPr.rFonts.set(
        qn("w:eastAsia"),
        "方正小标宋简体"
    )

    title_style.font.name = (
        "Times New Roman"
    )

    title_style.font.size = Pt(22)

    # ==========================
    # 页面设置
    # ==========================

    section = doc.sections[0]

    section.top_margin = Cm(2.5)

    section.bottom_margin = Cm(2.5)

    section.left_margin = Cm(3)

    section.right_margin = Cm(3)
# =====================================================
# Word生成模块
# =====================================================


def create_word(
        slides: list[SlideContent],
        output_file: Path
):
    """
    根据PPT内容生成Word
    """

    doc = Document()

    set_word_style(
        doc
    )

    # 文档标题

    title = doc.add_heading(
        "PPT转换报告",
        level=0
    )

    title.alignment = (
        WD_ALIGN_PARAGRAPH.CENTER
    )

    for slide in slides:

        # 添加页码提示

        doc.add_paragraph(
            f"[第 {slide.page} 页]"
        )

        # 添加标题

        doc.add_heading(
            slide.title,
            level=1
        )

        # 添加正文

        for paragraph in slide.paragraphs:

            p = doc.add_paragraph(
                paragraph
            )

            p.paragraph_format.space_after = Pt(8)

    doc.save(
        output_file
    )


# =====================================================
# 主程序
# =====================================================

def main():

    print(
        "开始读取PPT..."
    )

    if not PPT_FILE.exists():

        print(
            f"错误: 找不到文件 {PPT_FILE}"
        )

        return

    # 1.读取PPT

    slides = extract_ppt_text(
        PPT_FILE
    )

    print(
        f"共读取 {len(slides)} 页"
    )

    # 2.分析结构

    slides = analyze_content(
        slides
    )

    # 3.生成Word

    create_word(
        slides,
        WORD_FILE
    )

    print(
        "转换完成:"
    )

    print(
        WORD_FILE
    )


# =====================================================
# 程序入口
# =====================================================
if __name__ == "__main__":

    main()
    input("\n按 Enter 键退出...")
